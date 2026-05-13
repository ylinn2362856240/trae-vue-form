"""
分析用户手动修改后的 PPT 前三页的结构
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

PPT_PATH = r"D:\yangl-code\ai-form-demo\doc\KY_Week1_v2.pptx"
prs = Presentation(PPT_PATH)

print(f"Total slides: {len(prs.slides)}")

def dump(shape, indent=0):
    prefix = "  " * indent
    tp = str(shape.shape_type).split('(')[0].strip() if '(' in str(shape.shape_type) else str(shape.shape_type)
    pos = f"L={shape.left/914400:.2f} T={shape.top/914400:.2f}"
    size = f"W={shape.width/914400:.2f} H={shape.height/914400:.2f}"
    txt = ""
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                fs = ""
                if p.runs:
                    r = p.runs[0]
                    parts = []
                    if r.font.size: parts.append(f"sz={r.font.size/12700:.0f}pt")
                    if r.font.bold: parts.append("bold")
                    try:
                        if r.font.color and r.font.color.rgb: parts.append(f"c={r.font.color.rgb}")
                    except: pass
                    fs = f" ({', '.join(parts)})" if parts else ""
                txt += f"\n{prefix}    >> '{t[:70]}'{fs}"
    print(f"{prefix}[{tp:>10}] {pos}  {size}{txt}")
    if shape.shape_type == 6:
        try:
            for c in shape.shapes:
                dump(c, indent + 1)
        except: pass

for i in range(len(prs.slides)):
    slide = prs.slides[i]
    print(f"\n{'='*70}")
    print(f"SLIDE {i} (layout='{slide.slide_layout.name}')")
    print(f"{'='*70}")
    for shape in slide.shapes:
        dump(shape)
