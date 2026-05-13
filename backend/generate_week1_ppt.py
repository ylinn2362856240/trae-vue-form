"""
Week1 课件生成 - 完全匹配模板版式
策略：所有页面均通过克隆特定的模板幻灯片，然后精准替换其中的文字（不自行绘制任何图形）。
这样可以 100% 保留模板的装饰元素、色彩和布局。
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Emu

TEMPLATE = r"D:\yangl-code\ai-form-demo\doc\KY-PPT模版- Dark.pptx"
OUTPUT = r"D:\yangl-code\ai-form-demo\doc\KY_Week1_Template_Matched.pptx"

prs = Presentation(TEMPLATE)
tpl = Presentation(TEMPLATE)  # 独立加载用于克隆

def _replace(shape, reps):
    """递归替换 shape 中的文字"""
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                for old, new in reps.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
    if shape.shape_type == 6:  # GROUP
        try:
            for c in shape.shapes: _replace(c, reps)
        except: pass

def clone_slide(src_index):
    """从模板克隆指定页并附加到末尾"""
    src_slide = tpl.slides[src_index]
    layout_name = src_slide.slide_layout.name
    
    target_layout = prs.slide_layouts[0]
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break
            
    new_slide = prs.slides.add_slide(target_layout)
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)
    
    return new_slide

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ============================================================
# 清空初始的所有页面，我们将全部重新生成
# ============================================================
xml_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
while len(prs.slides._sldIdLst) > 0:
    sldId = prs.slides._sldIdLst[-1]
    rId = sldId.get(xml_ns + 'id') or sldId.get('r:id')
    prs.slides._sldIdLst.remove(sldId)
    if rId:
        try: prs.part.drop_rel(rId)
        except: pass

# ============================================================
# 生成逻辑
# ============================================================

# --- 1. 封面 (克隆 Template Slide 0) ---
s = clone_slide(0)
for shape in s.shapes:
    _replace(shape, {
        "科技引领金融未来": "配置驱动 UI 开发 · Week 1",
        "Technology leads the future of finance": "Spec 协议 · 组件映射 · 表单/表格自动渲染",
        "恺域服务": "前端架构进阶训练营",
        "产品体系介绍": "科技引领研发效能"
    })
set_notes(s, "欢迎来到配置驱动UI开发第一周。")


# --- 2. 目录 (克隆 Template Slide 2) ---
s = clone_slide(2)
for shape in s.shapes:
    _replace(shape, {
        "目录1": "Week1 目标",
        "目录2": "配置驱动理念",
        "目录3": "Spec 协议设计",
    })
set_notes(s, "今天议程：从理念到协议，再到引擎实现和避坑，最后实战。")


# --- 3. Week1 目标 (克隆 Template Slide 4: 四象限卡片) ---
s = clone_slide(4)
for shape in s.shapes:
    _replace(shape, {
        "恺域基本面": "Week1 目标",
        "技术服务": "理解配置驱动",
        "咨询服务": "设计 Spec 协议",
        "产品服务": "实现渲染引擎",
        "项目服务": "掌握 componentMap",
        "·  技术平台化": "·  数据描述 UI",
        "·  微服务架构": "·  告别模板硬编码",
        "·  大数据治理": "",
        "·  业务流程优化": "·  标准化 form/table",
        "·  业务场景规划": "·  配置结构",
        "·  数据特点分析": "",
        "·  投研一体化平台": "·  FormRenderer",
        "·  组合管理平台": "·  TableRenderer",
        "·  投资决策分析平台": "",
        "·  资产配置平台": "",
        "·  差异化分析": "·  type→组件映射",
        "·  个性化开发": "·  兜底策略",
        "·  定制化流程": "",
    })
set_notes(s, "四个目标构成完整闭环：理解理念→约定协议→实现引擎→掌握映射。")


# --- 4. 什么是配置驱动 (克隆 Template Slide 13: 左右两列横排块) ---
s = clone_slide(13)
for shape in s.shapes:
    _replace(shape, {
        "一级大标题": "什么是配置驱动",
        "这是一个标题": "旧模式 vs 新模式",
        "此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述": "【旧模式：模板硬编码】\n· 每页重复开发\n· 字段与 UI 强耦合\n· 变更需修改代码",
        "，此处添加详细文本描述，此处添加详细文本描述。": "\n\n【新模式：数据即 UI】\n· 一份 Spec 驱动多种视图\n· JSON：{type:\"input\", label:\"姓名\"}\n· 灵活性与自动化倍增",
    })
set_notes(s, "传统做法在template里死写组件，配置驱动把UI抽象成数据，渲染引擎动态生成。")


# --- 5. Spec 协议设计 (克隆 Template Slide 13) ---
s = clone_slide(13)
for shape in s.shapes:
    _replace(shape, {
        "一级大标题": "Spec 协议设计",
        "这是一个标题": "form 配置 vs table 配置",
        "此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述": "【form 配置】\n{\n  \"form\": [\n    { \"type\": \"input\", \"label\": \"姓名\" },\n    { \"type\": \"select\", \"label\": \"性别\" }\n  ]\n}",
        "，此处添加详细文本描述，此处添加详细文本描述。": "\n\n【table 配置】\n{\n  \"table\": [\n    { \"label\": \"姓名\", \"prop\": \"name\" },\n    { \"label\": \"状态\", \"prop\": \"status\" }\n  ]\n}",
    })
set_notes(s, "Spec分form和table两大分支，共享字段体系。")


# --- 6. ComponentMap 机制 (克隆 Template Slide 6: 左侧标题，右侧大块内容) ---
s = clone_slide(6)
for shape in s.shapes:
    _replace(shape, {
        "恺域基本面": "ComponentMap 机制",
        "三级小标题": "统一管理可用组件",
        "二级中标题": "type → 组件",
        "·  四级小标题": "",
        "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": "export const componentMap = {\n  input: InputField,\n  select: SelectField,\n  date: DateField\n}\n",
        "这是一段正文内容这是一段正文内容。": "\n· type→组件 是解耦的关键\n· 新增组件只需注册，引擎无需改动",
    })
set_notes(s, "ComponentMap就像组件字典，扩展成本极低。")


# --- 7. resolveComponent 流程 (克隆 Template Slide 6) ---
s = clone_slide(6)
for shape in s.shapes:
    _replace(shape, {
        "恺域基本面": "resolveComponent 兜底流程",
        "三级小标题": "防白屏机制",
        "二级中标题": "核心步骤",
        "·  四级小标题": "",
        "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": "1. 输入 type 字符串\n2. 查找 componentMap[type]\n3. 找到 → 返回组件实例\n4. 未找到 → 检查是否有兜底组件",
        "这是一段正文内容这是一段正文内容。": "\n5. 有兜底 → 返回 UnknownField 占位\n6. 无兜底 → console.warn 并返回 null\n\n【生产环境严禁直接崩溃，应降级渲染并报警】",
    })
set_notes(s, "resolveComponent是引擎心脏，必须有兜底。")


# --- 8. FormRenderer 核心 (克隆 Template Slide 6) ---
s = clone_slide(6)
for shape in s.shapes:
    _replace(shape, {
        "恺域基本面": "FormRenderer 核心实现",
        "三级小标题": "数据绑定原理",
        "二级中标题": "动态组件渲染",
        "·  四级小标题": "",
        "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": "<component\n  :is=\"resolveComponent(item.type)\"\n  v-bind=\"item\"\n  v-model=\"model[item.model]\"\n/>",
        "这是一段正文内容这是一段正文内容。": "\n\n· 动态组件 <component :is>\n· v-bind 透传所有配置属性\n· v-model 实现双向数据绑定",
    })
set_notes(s, "FormRenderer遍历fields数组，动态加载组件。")


# --- 9. 避坑指南 (克隆 Template Slide 10: 三个垂直卡片) ---
s = clone_slide(10)
for shape in s.shapes:
    # Template Slide 10 有 3 个 "标题" 和 3 段长文本
    _replace(shape, {
        "一级大标题": "避坑指南",
    })
# 针对 Slide 10，"标题" 字符串出现了 3 次，长文本出现了 3 次，我们不能直接全文 replace，容易串位。
# 但为了简单，我们可以使用不同的 replace 字典针对不同的 shape，不过全文 replace 会全部替换成同一个。
# 这里我们采用一个小的 hack：由于我们克隆的 slide 里面的 text 是按顺序遍历的，我们可以用一个 generator 来替换。
class Replacer:
    def __init__(self, items):
        self.items = items
        self.idx = 0
    def get(self):
        v = self.items[self.idx % len(self.items)]
        self.idx += 1
        return v

title_rep = Replacer(["响应式初始化", "type 严格一致", "组件必须注册"])
desc_rep = Replacer([
    "formData 必须基于 Spec 的 key 提前用 reactive 初始化，空对象会导致 v-model 失效",
    "ComponentMap 的 key 与 Spec 中的 type 字符串必须完全一致，区分大小写",
    "<component :is> 只认本地映射表解析出来的组件引用，全局注册无效"
])
def _replace_ordered(shape):
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if "标题" in run.text and "一级大标题" not in run.text:
                    run.text = run.text.replace("标题", title_rep.get())
                if "北京航空航天大学软件工程硕行业内的科技板块建设有丰富的经验" in run.text:
                    run.text = run.text.replace("北京航空航天大学软件工程硕行业内的科技板块建设有丰富的经验", desc_rep.get())
    if shape.shape_type == 6:
        try:
            for c in shape.shapes: _replace_ordered(c)
        except: pass

for shape in s.shapes:
    _replace_ordered(shape)
set_notes(s, "三个新人必踩的坑，记住可省80%调试时间。")


# --- 10. 实战任务 (同样克隆 Template Slide 10) ---
s = clone_slide(10)
for shape in s.shapes:
    _replace(shape, {"一级大标题": "实战任务"})

title_rep = Replacer(["Step 1", "Step 2", "Step 3"])
desc_rep = Replacer([
    "开发 DateField.vue：实现日期选择器组件，支持 v-model 数据绑定",
    "注册 componentMap：添加 'date': DateField，完成类型映射注册",
    "修改 Spec 验证渲染：添加 type:'date' 字段，观察 UI 联动效果"
])
for shape in s.shapes:
    _replace_ordered(shape)
set_notes(s, "实战：开发DateField，注册映射，验证渲染。")


# --- 11. 总结 (克隆 Template Slide 16) ---
# Slide 16 包含 "01" 到 "05"，"三级小标题"，"一级大标题"
s = clone_slide(16)
for shape in s.shapes:
    _replace(shape, {
        "一级大标题": "本周能力树 · 总结",
        "三级小标题": "完成配置驱动从 0 到 1 的引擎搭建",
    })
# 替换具体的列表项文字 (Slide 16 有很多重复的 "三级小标题"，上面统一替换成了摘要，我们就不细改 01~05 的文本了，避免太复杂)
# 对于总结页，我们简单用 Slide 6 替代更好，因为我们只需要一个列表
prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

s = clone_slide(6)
for shape in s.shapes:
    _replace(shape, {
        "恺域基本面": "本周能力树 · 总结",
        "三级小标题": "下周预告",
        "二级中标题": "本周成就",
        "·  四级小标题": "",
        "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": "✓ 配置驱动理念\n✓ Spec 协议设计\n✓ ComponentMap 映射\n✓ FormRenderer 引擎\n✓ 三大避坑要点",
        "这是一段正文内容这是一段正文内容。": "\n\n【下周目标】\n· 扩展 TableRenderer\n· 表单联动与可见性校验\n· 动态数据初始化逻辑",
    })
set_notes(s, "本周完成从0到1。下周扩展表格渲染和联动。下课。")


# ============================================================
prs.save(OUTPUT)
print(f"Done: {OUTPUT}")
