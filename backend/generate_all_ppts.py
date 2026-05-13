"""
批量生成 Week2 到 Week5 课件
完全按照 KY_Week1_v2 的模板映射策略
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

TEMPLATE = r"D:\yangl-code\ai-form-demo\doc\KY-PPT模版- Dark.pptx"

def _replace(shape, reps):
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                for old, new in reps.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
    if shape.shape_type == 6:
        try:
            for c in shape.shapes: _replace(c, reps)
        except: pass

class Replacer:
    def __init__(self, items):
        self.items = items
        self.idx = 0
    def get(self):
        v = self.items[self.idx % len(self.items)]
        self.idx += 1
        return v

def _replace_ordered(shape, title_rep, desc_rep):
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if "标题" in run.text and "一级大标题" not in run.text:
                    run.text = run.text.replace("标题", title_rep.get())
                if "北京航空航天大学软件工程硕行业内的科技板块建设有丰富的经验" in run.text:
                    run.text = run.text.replace("北京航空航天大学软件工程硕行业内的科技板块建设有丰富的经验", desc_rep.get())
    if shape.shape_type == 6:
        try:
            for c in shape.shapes: _replace_ordered(c, title_rep, desc_rep)
        except: pass

def generate_ppt(week_data, output_path):
    prs = Presentation(TEMPLATE)
    tpl = Presentation(TEMPLATE)
    
    def clone_slide(src_idx):
        src_slide = tpl.slides[src_idx]
        target_layout = prs.slide_layouts[0]
        for layout in prs.slide_layouts:
            if layout.name == src_slide.slide_layout.name:
                target_layout = layout
                break
        new_slide = prs.slides.add_slide(target_layout)
        for shape in src_slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape.element))
        return new_slide
        
    # Clear init slides
    xml_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    while len(prs.slides._sldIdLst) > 0:
        sldId = prs.slides._sldIdLst[-1]
        rId = sldId.get(xml_ns + 'id') or sldId.get('r:id')
        prs.slides._sldIdLst.remove(sldId)
        if rId:
            try: prs.part.drop_rel(rId)
            except: pass

    # 1. Cover
    s = clone_slide(0)
    for shape in s.shapes:
        _replace(shape, {
            "科技引领金融未来": week_data["cover_title"],
            "Technology leads the future of finance": week_data["cover_sub"],
            "恺域服务": "前端架构进阶训练营",
            "产品体系介绍": "科技引领研发效能"
        })

    # 2. TOC
    s = clone_slide(2)
    for shape in s.shapes:
        _replace(shape, {
            "目录1": week_data["toc_1"],
            "目录2": week_data["toc_2"],
            "目录3": week_data["toc_3"],
        })

    # 3. Quadrants
    s = clone_slide(4)
    q = week_data["quadrants"]
    for shape in s.shapes:
        _replace(shape, {
            "恺域基本面": week_data["toc_1"],
            "技术服务": q["q1_title"],
            "咨询服务": q["q2_title"],
            "产品服务": q["q3_title"],
            "项目服务": q["q4_title"],
            "·  技术平台化": q["q1_l1"],
            "·  微服务架构": q["q1_l2"],
            "·  大数据治理": "",
            "·  业务流程优化": q["q2_l1"],
            "·  业务场景规划": q["q2_l2"],
            "·  数据特点分析": "",
            "·  投研一体化平台": q["q3_l1"],
            "·  组合管理平台": q["q3_l2"],
            "·  投资决策分析平台": "",
            "·  资产配置平台": "",
            "·  差异化分析": q["q4_l1"],
            "·  个性化开发": q["q4_l2"],
            "·  定制化流程": "",
        })

    # 4. Section 1 (2 Columns)
    s = clone_slide(13)
    s1 = week_data["sec1"]
    for shape in s.shapes:
        _replace(shape, {
            "一级大标题": s1["title"],
            "这是一个标题": s1["sub"],
            "此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述": s1["col1"],
            "，此处添加详细文本描述，此处添加详细文本描述。": s1["col2"],
        })

    # 5. Section 2 (1 block / code)
    s = clone_slide(6)
    s2 = week_data["sec2"]
    for shape in s.shapes:
        _replace(shape, {
            "恺域基本面": s2["title"],
            "三级小标题": s2["sub1"],
            "二级中标题": s2["sub2"],
            "·  四级小标题": "",
            "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": s2["code"],
            "这是一段正文内容这是一段正文内容。": s2["desc"],
        })

    # 6. Section 3 (2 Columns or another block)
    s = clone_slide(13)
    s3 = week_data["sec3"]
    for shape in s.shapes:
        _replace(shape, {
            "一级大标题": s3["title"],
            "这是一个标题": s3["sub"],
            "此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述，此处添加详细文本描述": s3["col1"],
            "，此处添加详细文本描述，此处添加详细文本描述。": s3["col2"],
        })

    # 7. Warnings
    s = clone_slide(10)
    for shape in s.shapes:
        _replace(shape, {"一级大标题": "避坑指南"})
    title_rep = Replacer(week_data["warn_titles"])
    desc_rep = Replacer(week_data["warn_descs"])
    for shape in s.shapes:
        _replace_ordered(shape, title_rep, desc_rep)

    # 8. Tasks
    s = clone_slide(10)
    for shape in s.shapes:
        _replace(shape, {"一级大标题": "实战任务"})
    title_rep = Replacer(["Step 1", "Step 2", "Step 3"])
    desc_rep = Replacer(week_data["tasks"])
    for shape in s.shapes:
        _replace_ordered(shape, title_rep, desc_rep)

    # 9. Summary
    s = clone_slide(6)
    for shape in s.shapes:
        _replace(shape, {
            "恺域基本面": "本周能力树 · 总结",
            "三级小标题": "下周预告",
            "二级中标题": "本周成就",
            "·  四级小标题": "",
            "这是一段正文内容这是一段正文内容这是一段正文内容这是一段正文内容": week_data["summary"],
            "这是一段正文内容这是一段正文内容。": week_data["next_week"],
        })

    prs.save(output_path)
    print(f"Generated: {output_path}")

# ================= Data Definitions =================
weeks = [
    {
        "week": 2,
        "cover_title": "TableRenderer 与表单联动 · Week 2",
        "cover_sub": "复杂表格渲染 · 表达式引擎 · 动态联动逻辑",
        "toc_1": "Week2 目标", "toc_2": "TableRenderer", "toc_3": "动态联动",
        "quadrants": {
            "q1_title": "理解 Table", "q1_l1": "·  开发 TableRenderer", "q1_l2": "·  支持插槽与操作列",
            "q2_title": "设计联动协议", "q2_l1": "·  v-if 表达式", "q2_l2": "·  值联动公式",
            "q3_title": "表达式引擎", "q3_l1": "·  安全执行 eval", "q3_l2": "·  沙箱隔离机制",
            "q4_title": "实战落地", "q4_l1": "·  完成表格展示", "q4_l2": "·  表单联动实战"
        },
        "sec1": {
            "title": "TableRenderer 设计思路", "sub": "常规表格开发 vs 配置驱动表格",
            "col1": "【常规表格】\n· el-table-column 硬编码\n· 业务逻辑与 UI 耦合\n· 增加字段需改代码",
            "col2": "\n\n【配置驱动表格】\n· 一份 JSON 描述列信息\n· 动态解析 prop 和 label\n· 自定义列通过 slot 动态注入"
        },
        "sec2": {
            "title": "TableRenderer 核心实现", "sub1": "动态列渲染", "sub2": "核心组件解析",
            "code": "<el-table :data=\"data\">\n  <template v-for=\"col in schema\">\n    <el-table-column v-bind=\"col\">\n      <template #default=\"{row}\" v-if=\"col.slot\">\n        <slot :name=\"col.slot\" :row=\"row\" />\n      </template>\n    </el-table-column>\n  </template>\n</el-table>",
            "desc": "\n· 遍历 schema 动态生成列\n· v-bind 透传属性\n· 具名插槽支持高度定制（如操作按钮）"
        },
        "sec3": {
            "title": "表单联动协议设计", "sub": "显隐联动 vs 值联动",
            "col1": "【显隐联动 (Visibility)】\n· 控制组件 v-if/v-show\n· 协议：{ showIf: 'model.age > 18' }\n· 满足条件才渲染该表单项",
            "col2": "\n\n【值联动 (Value Linkage)】\n· 当 A 改变时自动计算 B\n· 协议：{ valueExpr: 'model.price * qty' }\n· 实时计算响应式更新"
        },
        "warn_titles": ["坑1 · XSS 风险", "坑2 · 循环依赖", "坑3 · 插槽作用域"],
        "warn_descs": [
            "严禁用联动表达式执行外部不可信输入，必须限制在 model 属性范围内",
            "A 的显隐依赖 B，B 的值又依赖 A，会导致无限更新循环，需做防御",
            "Table 的 slot 必须正确传递 row 和 index，否则外层无法获取当前行数据"
        ],
        "tasks": [
            "开发 TableRenderer，实现基础配置列和操作列 slot",
            "在 FormRenderer 中集成 evaluateExpr，实现 showIf 显隐联动",
            "编写测试 Spec，验证表格渲染和联动表单"
        ],
        "summary": "✓ TableRenderer 实现\n✓ 动态列与插槽机制\n✓ 联动协议设计\n✓ 安全表达式引擎\n✓ 循环依赖防御",
        "next_week": "\n\n【下周预告】\n· 复杂业务组件封装\n· 状态管理架构优化\n· 跨组件通信方案"
    },
    {
        "week": 3,
        "cover_title": "复杂组件与状态管理 · Week 3",
        "cover_sub": "Upload/Tree 封装 · Pinia 状态树 · 跨层通信",
        "toc_1": "Week3 目标", "toc_2": "复杂组件", "toc_3": "状态管理",
        "quadrants": {
            "q1_title": "复杂组件", "q1_l1": "·  封装 Upload/Tree", "q1_l2": "·  统一 v-model 接口",
            "q2_title": "状态管理", "q2_l1": "·  引入 Pinia", "q2_l2": "·  管理全局状态",
            "q3_title": "跨层通信", "q3_l1": "·  Provide / Inject", "q3_l2": "·  解决 Prop 钻取",
            "q4_title": "工程规范", "q4_l1": "·  统一事件总线", "q4_l2": "·  标准化数据流"
        },
        "sec1": {
            "title": "复杂组件封装思路", "sub": "原子组件 vs 复合组件",
            "col1": "【原子组件 (Input/Select)】\n· 数据结构简单 (String/Number)\n· 无内部状态\n· 直接透传 v-model",
            "col2": "\n\n【复合组件 (Upload/Tree)】\n· 数据结构复杂 (Array/Object)\n· 有内部请求与交互状态\n· 需要包装适配统一接口"
        },
        "sec2": {
            "title": "Upload 组件适配层", "sub1": "统一接口", "sub2": "数据转换",
            "code": "<template>\n  <el-upload\n    :file-list=\"internalList\"\n    @change=\"handleChange\"\n  />\n</template>\n<script>\n// 将 el-upload 格式转为标准 Array<URL>\nconst emit = defineEmits(['update:modelValue']);\n</script>",
            "desc": "\n· 隐藏组件内部实现细节\n· 对外只暴露标准 modelValue\n· 在 componentMap 中无缝注册"
        },
        "sec3": {
            "title": "全局状态架构", "sub": "Pinia vs Provide",
            "col1": "【Pinia Store】\n· 抽离 UI 层与逻辑层\n· Spec 存储供多组件共享\n· 便于做撤销/重做功能",
            "col2": "\n\n【Provide / Inject】\n· 解决深度嵌套组件传参问题\n· 将 Engine Context 注入到底层 Field\n· 保持组件签名的干净"
        },
        "warn_titles": ["坑1 · 异步陷阱", "坑2 · Store 污染", "坑3 · 深层监听"],
        "warn_descs": [
            "Upload 组件状态是异步的，提交表单前必须等待所有文件上传完毕",
            "多个 Form 实例同时渲染时，Pinia Store 必须区分 namespace",
            "对象类型的值在复合组件内变更时，可能触发不了外层的 watch"
        ],
        "tasks": [
            "封装 UploadField 组件，屏蔽内部 FileList 细节",
            "引入 Pinia，创建 useEngineStore 管理全局 Spec",
            "实现基于 Provide/Inject 的跨层级联动上下文"
        ],
        "summary": "✓ 复合组件适配封装\n✓ 异步状态同步处理\n✓ Pinia 状态树架构\n✓ Context 跨层通信\n✓ 多实例命名空间隔离",
        "next_week": "\n\n【下周预告】\n· FastAPI 后端架构设计\n· AI 服务接入\n· 稳定性兜底策略"
    },
    {
        "week": 4,
        "cover_title": "后端服务与兜底逻辑 · Week 4",
        "cover_sub": "FastAPI 架构 · Mock 降级策略 · 高可用设计",
        "toc_1": "Week4 目标", "toc_2": "服务端架构", "toc_3": "兜底策略",
        "quadrants": {
            "q1_title": "FastAPI 架构", "q1_l1": "·  构建高性能后端", "q1_l2": "·  定义标准接口",
            "q2_title": "LLM 接入", "q2_l1": "·  对接大模型 API", "q2_l2": "·  处理流式请求",
            "q3_title": "Mock 降级", "q3_l1": "·  无 AI 时兜底", "q3_l2": "·  规则引擎生成 UI",
            "q4_title": "稳定性保障", "q4_l1": "·  超时与防抖处理", "q4_l2": "·  异常捕获与重试"
        },
        "sec1": {
            "title": "系统高可用架构", "sub": "直接调用 vs 服务端代理",
            "col1": "【前端直连大模型】\n· 密钥暴露风险极高\n· 无法做请求限流管控\n· 跨域问题难以优雅解决",
            "col2": "\n\n【服务端代理层 (API)】\n· 保护 API Key 等敏感信息\n· 统一接口，方便切换模型引擎\n· 支持增加缓存和 Mock 策略"
        },
        "sec2": {
            "title": "FastAPI 接口设计", "sub1": "/generate 路由", "sub2": "Pydantic 校验",
            "code": "@app.post(\"/generate\")\nasync def generate_ui(req: GenerateRequest):\n    try:\n        spec = await llm.generate_spec(req.prompt)\n        return {\"code\": 0, \"data\": spec}\n    except Exception as e:\n        return mock_generate(req.prompt)",
            "desc": "\n· 定义强类型的请求体\n· 主干逻辑尝试调用 AI\n· 异常时静默降级为 Mock 数据"
        },
        "sec3": {
            "title": "Mock 降级策略", "sub": "硬编码 vs 规则引擎",
            "col1": "【硬编码兜底】\n· 只返回固定的一套表单\n· 用户体验割裂\n· 无法响应业务提示词",
            "col2": "\n\n【规则引擎降级】\n· 根据 Prompt 关键词正则匹配\n· 动态拼装预设字段 (如手机号)\n· 体验平滑，用户感知不到异常"
        },
        "warn_titles": ["坑1 · 超时阻塞", "坑2 · JSON 解析", "坑3 · 跨域问题"],
        "warn_descs": [
            "LLM 请求往往在 10s 以上，必须设置 Timeout 并尽早降级，防止前端假死",
            "大模型返回的未必是纯 JSON，可能有 Markdown 代码块包裹，必须清洗",
            "前后端分离部署时，FastAPI 必须正确配置 CORS 允许前端端口访问"
        ],
        "tasks": [
            "使用 FastAPI 搭建后端服务，配置跨域",
            "实现大模型请求服务，并编写 Mock 降级引擎",
            "前后端联调，模拟断网测试降级体验"
        ],
        "summary": "✓ FastAPI 后端搭建\n✓ 大模型接口集成\n✓ 智能规则 Mock 引擎\n✓ JSON 字符串清洗\n✓ 全局异常兜底架构",
        "next_week": "\n\n【下周预告】\n· 提示词结构化设计\n· Few-Shot 实战\n· 完整项目闭环验收"
    },
    {
        "week": 5,
        "cover_title": "AI 提示词工程与实战 · Week 5",
        "cover_sub": "Prompt 设计 · 结构约束 · 智能生成闭环",
        "toc_1": "Week5 目标", "toc_2": "Prompt 工程", "toc_3": "系统集成",
        "quadrants": {
            "q1_title": "提示词原理", "q1_l1": "·  掌握角色设定", "q1_l2": "·  输出约束条件",
            "q2_title": "结构化输出", "q2_l1": "·  Few-Shot 示例", "q2_l2": "·  严控 JSON 格式",
            "q3_title": "容错重试", "q3_l1": "·  解析失败修复", "q3_l2": "·  自动重试机制",
            "q4_title": "闭环验收", "q4_l1": "·  前后端全栈打通", "q4_l2": "·  一句话生成页面"
        },
        "sec1": {
            "title": "大模型引导策略", "sub": "开放式 vs 约束式 Prompt",
            "col1": "【开放式 Prompt】\n· \"帮我生成一个登录表单\"\n· 每次返回的格式都不一样\n· 无法直接送入 UI 引擎渲染",
            "col2": "\n\n【约束式 Prompt (System)】\n· \"你是前端专家，严格按 JSON 输出...\"\n· 提供范例 (Few-Shot)\n· 输出具备高度确定性"
        },
        "sec2": {
            "title": "高质量 Prompt 模板", "sub1": "系统指令", "sub2": "约束示范",
            "code": "SYSTEM_PROMPT = \"\"\"\n你是一个表单引擎配置生成器。\n必须返回合法的 JSON，不要输出任何解释。\n支持的 type: [input, number, date, select]\n示例格式:\n{ \"form\": [{ \"type\": \"input\", ... }] }\n\"\"\"",
            "desc": "\n· 明确身份定位\n· 严格限定组件库可用范围\n· 提供标准的 JSON 骨架"
        },
        "sec3": {
            "title": "JSON 解析容错机制", "sub": "单纯 loads vs 修复正则",
            "col1": "【直接 json.loads】\n· 极易因为尾部逗号报错\n· 无法处理 Markdown 标记\n· 成功率通常低于 80%",
            "col2": "\n\n【清洗与修复 (Clean+Regex)】\n· 剥离 ```json 与 ``` 标签\n· 正则处理多余标点\n· 失败时让 LLM 重新 \"修复这个 JSON\""
        },
        "warn_titles": ["坑1 · 幻觉组件", "坑2 · 模型记忆", "坑3 · 提示词过长"],
        "warn_descs": [
            "大模型极易凭空创造 `type: 'map'` 等不存在的组件，必须白名单限制",
            "单轮对话无状态，需将上次生成的错误 Spec 喂回给模型要求修改",
            "超长示例会消耗大量 Token 且容易失焦，应提炼最简 Few-Shot"
        ],
        "tasks": [
            "设计 System Prompt，确保 95% 以上 JSON 格式正确率",
            "在后端加入 Markdown 清洗逻辑和容错解析机制",
            "完成最终闭环，演示对话生成多种业务场景的表单"
        ],
        "summary": "✓ 结构化 Prompt 设计\n✓ Few-Shot 示例调优\n✓ JSON 鲁棒性清洗\n✓ 全栈应用贯通联调\n✓ 幻觉控制与校验",
        "next_week": "\n\n【完结撒花】\n恭喜完成《配置驱动 UI 开发》全部课程！\n科技引领研发效能！"
    }
]

for w in weeks:
    generate_ppt(w, f"D:\\yangl-code\\ai-form-demo\\doc\\KY_Week{w['week']}_PPT.pptx")
