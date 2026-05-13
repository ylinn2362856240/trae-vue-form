import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
TEMPLATE = r"D:\yangl-code\ai-form-demo\doc\KY-PPT模版- Dark.pptx"
prs = Presentation(TEMPLATE)

def dump_slide(idx):
    print(f"--- Slide {idx} ---")
    for shape in prs.slides[idx].shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                print(f"[{shape.name}] text: '{p.text}'")
        if shape.shape_type == 6:
            try:
                for c in shape.shapes:
                    if c.has_text_frame:
                        for p in c.text_frame.paragraphs:
                            print(f"[{c.name}] text: '{p.text}'")
                    if c.shape_type == 6:
                        try:
                            for cc in c.shapes:
                                if cc.has_text_frame:
                                    for p in cc.text_frame.paragraphs:
                                        print(f"[{cc.name}] text: '{p.text}'")
                        except: pass
            except: pass

dump_slide(10)
dump_slide(13)
