"""
Extract text from specific template slides for mapping replacement
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
TEMPLATE = r"D:\yangl-code\ai-form-demo\doc\KY-PPT模版- Dark.pptx"
prs = Presentation(TEMPLATE)

def extract_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip(): texts.append(p.text.strip())
        if shape.shape_type == 6:
            try:
                for c in shape.shapes:
                    if hasattr(c, 'text_frame') and c.has_text_frame:
                        for p in c.text_frame.paragraphs:
                            if p.text.strip(): texts.append(p.text.strip())
            except: pass
    return texts

print("Slide 6 (1 column?):", extract_text(prs.slides[6]))
print("Slide 7 (4 vertical steps?):", extract_text(prs.slides[7]))
print("Slide 10 (3 items?):", extract_text(prs.slides[10]))
print("Slide 11 (4 columns?):", extract_text(prs.slides[11]))
print("Slide 13 (2 items?):", extract_text(prs.slides[13]))
print("Slide 16 (3 items?):", extract_text(prs.slides[16]))
